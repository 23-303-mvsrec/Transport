from pptx import Presentation

def inspect_details(pptx_path):
    prs = Presentation(pptx_path)
    for i, slide in enumerate(prs.slides, 1):
        print(f"=== SLIDE {i} ===")
        print(f"Title: {slide.shapes.title.text if slide.shapes.title else '[No Title]'}")
        
        # Check notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                print(f"Notes:\n{notes}")
        
        # Check all shape types and properties
        print("Shapes:")
        for j, shape in enumerate(slide.shapes):
            shape_type = shape.shape_type
            name = shape.name
            text = shape.text_frame.text.strip() if shape.has_text_frame else ""
            print(f"  - Shape {j}: Type={shape_type}, Name='{name}', Text Length={len(text)}")
            if text:
                print(f"    Text: {text[:100]}...")
        print("\n" + "="*40 + "\n")

if __name__ == "__main__":
    inspect_details("D24_ProjectReview2 (1).pptx")
