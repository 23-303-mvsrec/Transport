import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ----------------------------------------------------
    # COLOR PALETTE DESIGN
    # ----------------------------------------------------
    # Dark Slate (Intro/Outro)
    DARK_BG = RGBColor(10, 25, 47)
    DARK_TEXT_PRIMARY = RGBColor(255, 255, 255)
    DARK_TEXT_SECONDARY = RGBColor(136, 146, 176)
    ACCENT_TEAL = RGBColor(100, 255, 218)
    
    # Light Slate (Content slides)
    LIGHT_BG = RGBColor(247, 250, 252)
    LIGHT_TEXT_PRIMARY = RGBColor(26, 54, 93)      # Deep Navy
    LIGHT_TEXT_SECONDARY = RGBColor(74, 85, 104)   # Charcoal
    LIGHT_ACCENT_TEAL = RGBColor(49, 151, 149)     # Medium Teal
    LIGHT_ACCENT_ORANGE = RGBColor(221, 107, 32)   # Coral Accent
    
    blank_layout = prs.slide_layouts[6]
    
    # ----------------------------------------------------
    # HELPER FUNCTIONS FOR CLEAN DESIGN
    # ----------------------------------------------------
    
    def apply_solid_background(slide, color):
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = color
        bg_shape.line.color.rgb = color
        return bg_shape

    def add_slide_header(slide, title_text):
        # Header Text Box
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0)
        tf.margin_top = Inches(0)
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Segoe UI"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = LIGHT_TEXT_PRIMARY
        
        # Divider Line below title
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.25), Inches(11.733), Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = LIGHT_ACCENT_TEAL
        line.line.color.rgb = LIGHT_ACCENT_TEAL

    def configure_paragraph(p, text, size=15, bold=False, color=LIGHT_TEXT_SECONDARY, space_after=12, level=0):
        p.text = text
        p.font.name = "Segoe UI"
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(space_after)
        p.level = level
        p.alignment = PP_ALIGN.LEFT

    # ----------------------------------------------------
    # SLIDE 1: COVER SLIDE (DARK SLATE THEME)
    # ----------------------------------------------------
    slide_1 = prs.slides.add_slide(blank_layout)
    apply_solid_background(slide_1, DARK_BG)
    
    # Large College Header Text Box
    coll_box = slide_1.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(1.2))
    ctf = coll_box.text_frame
    ctf.word_wrap = True
    ctf.margin_left = Inches(0)
    
    p = ctf.paragraphs[0]
    configure_paragraph(p, "Matrusri Engineering College", size=18, bold=True, color=DARK_TEXT_PRIMARY, space_after=2)
    p = ctf.add_paragraph()
    configure_paragraph(p, "An Autonomous Institution | Approved by AICTE, Affiliated to Osmania University", size=11, color=DARK_TEXT_SECONDARY, space_after=2)
    p = ctf.add_paragraph()
    configure_paragraph(p, "Department of Computer Science and Engineering", size=12, bold=True, color=ACCENT_TEAL, space_after=0)

    # Main Project Title Box
    title_box = slide_1.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.733), Inches(2.2))
    ttf = title_box.text_frame
    ttf.word_wrap = True
    ttf.margin_left = Inches(0)
    
    p = ttf.paragraphs[0]
    configure_paragraph(p, "CITYBUS", size=48, bold=True, color=ACCENT_TEAL, space_after=8)
    p = ttf.add_paragraph()
    configure_paragraph(p, "Live Public Transport Tracking & Fleet Management System", size=22, bold=True, color=DARK_TEXT_PRIMARY, space_after=8)
    p = ttf.add_paragraph()
    configure_paragraph(p, "Real-Time Telemetry, Adaptive Mapping, and Progressive Web App Integration", size=14, color=DARK_TEXT_SECONDARY, space_after=0)
    
    # Presenters & Guidance Box (2 columns at bottom)
    info_box = slide_1.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(11.733), Inches(2.0))
    itf = info_box.text_frame
    itf.word_wrap = True
    itf.margin_left = Inches(0)
    
    # Left column in text frame via spaces/tabs or layout
    p = itf.paragraphs[0]
    configure_paragraph(p, "Under the Guidance of:                          Presented By:", size=13, bold=True, color=DARK_TEXT_PRIMARY, space_after=4)
    p = itf.add_paragraph()
    configure_paragraph(p, "Mrs. P. Uma Maheshwari                           [Your Name] ([Your Roll No])", size=13, color=DARK_TEXT_SECONDARY, space_after=2)
    p = itf.add_paragraph()
    configure_paragraph(p, "Assistant Professor                              Batch No: D24 | Project Review - II", size=13, color=DARK_TEXT_SECONDARY, space_after=0)

    # ----------------------------------------------------
    # SLIDE 2: PROBLEM DEFINITION & OBJECTIVE (2 COLUMNS)
    # ----------------------------------------------------
    slide_2 = prs.slides.add_slide(blank_layout)
    apply_solid_background(slide_2, LIGHT_BG)
    add_slide_header(slide_2, "Problem Definition & Objectives")
    
    # Left Box (Problem)
    left_box = slide_2.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    ltf = left_box.text_frame
    ltf.word_wrap = True
    ltf.margin_left = Inches(0)
    p = ltf.paragraphs[0]
    configure_paragraph(p, "Problem Definition", size=18, bold=True, color=LIGHT_ACCENT_TEAL, space_after=10)
    
    problems = [
        "Unpredictable Transit wait times: Commuters lack real-time bus locations, leading to massive schedules uncertainty.",
        "Manual telemetry & load logs: Operators track cabins capacity and duty timelines on paper sheets, causing delayed data updates.",
        "High proprietary GPS costs: Fleet authorities rely on expensive proprietary hardware, avoiding software solutions.",
        "Low network reliability: Web applications freeze or fail in low-bandwidth regions due to rigid mapping dependencies.",
        "Lack of Simulative Testbeds: Live route calculations can't be debugged locally without driving actual roads."
    ]
    for prob in problems:
        p = ltf.add_paragraph()
        configure_paragraph(p, prob, size=13, color=LIGHT_TEXT_SECONDARY, space_after=8, level=1)

    # Right Box (Objectives)
    right_box = slide_2.shapes.add_textbox(Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.2))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    rtf.margin_left = Inches(0)
    p = rtf.paragraphs[0]
    configure_paragraph(p, "Objectives", size=18, bold=True, color=LIGHT_ACCENT_TEAL, space_after=10)
    
    objs = [
        "Commuter Map Tracking: Deliver interactive route maps, bus markers, stop lists, and calculated coordinates-based ETAs.",
        "Mobile Driver Console: Build a duty console streaming speed, coordinates, and manual occupant counts in real time.",
        "Fleet Control Console: Provision an admin center mapping active vehicles and database configurations.",
        "Connection Adaptive Engine: Dynamically swap from Google Maps to Leaflet Maps under low-network alerts.",
        "Resilient Dual-Sync Mode: Sync to Firestore if online, or fall back to client-side physics loop simulations offline."
    ]
    for obj in objs:
        p = rtf.add_paragraph()
        configure_paragraph(p, obj, size=13, color=LIGHT_TEXT_SECONDARY, space_after=8, level=1)

    # ----------------------------------------------------
    # SLIDE 3: SYSTEM REQUIREMENTS (2 COLUMNS)
    # ----------------------------------------------------
    slide_3 = prs.slides.add_slide(blank_layout)
    apply_solid_background(slide_3, LIGHT_BG)
    add_slide_header(slide_3, "System Requirements")
    
    # Left Box (Software)
    left_box = slide_3.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    ltf = left_box.text_frame
    ltf.word_wrap = True
    ltf.margin_left = Inches(0)
    p = ltf.paragraphs[0]
    configure_paragraph(p, "Software Specifications", size=18, bold=True, color=LIGHT_ACCENT_TEAL, space_after=10)
    
    sw_specs = [
        "Frontend: React 18, Vite bundling tooling",
        "Styling: Tailwind CSS 3, PostCSS, Vanilla CSS3 variables",
        "State & Routing: React Router DOM 6, Custom Contexts",
        "Database (Cloud): Firebase Cloud Firestore (Real-time NoSQL)",
        "Authentication: Firebase Authentication",
        "PWA Support: Vite PWA Plugin, Service Workers, Workbox",
        "API Integrations: Google Maps JS API, Leaflet / OpenStreetMap",
        "Offline Database: Browser LocalStorage Client Cache",
        "Development Tools: VS Code, Git / GitHub, Windows 10/11"
    ]
    for spec in sw_specs:
        p = ltf.add_paragraph()
        configure_paragraph(p, spec, size=13, color=LIGHT_TEXT_SECONDARY, space_after=6, level=1)

    # Right Box (Hardware)
    right_box = slide_3.shapes.add_textbox(Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.2))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    rtf.margin_left = Inches(0)
    p = rtf.paragraphs[0]
    configure_paragraph(p, "Hardware Specifications", size=18, bold=True, color=LIGHT_ACCENT_TEAL, space_after=10)
    
    hw_specs = [
        "Development Processor: Intel Core i5/i7 (or equivalent) or higher",
        "Development RAM: 8 GB minimum (16 GB recommended)",
        "Development Storage: 256 GB SSD (minimum 5 GB free space)",
        "Internet Connection: Broadband for sync setup; 3G/4G/5G mobile data for field coordinates simulation",
        "End-User Device: Android smartphone, tablet, or iOS phone running a modern browser (Chrome, Safari, Firefox)",
        "Geolocation Node: Device built-in GPS/GNSS receiver module"
    ]
    for spec in hw_specs:
        p = rtf.add_paragraph()
        configure_paragraph(p, spec, size=13, color=LIGHT_TEXT_SECONDARY, space_after=8, level=1)

    # ----------------------------------------------------
    # SLIDE 4: EXISTING SYSTEM & LIMITATIONS (2 COLUMNS)
    # ----------------------------------------------------
    slide_4 = prs.slides.add_slide(blank_layout)
    apply_solid_background(slide_4, LIGHT_BG)
    add_slide_header(slide_4, "Existing System & Limitations")
    
    # Left Box (Existing)
    left_box = slide_4.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    ltf = left_box.text_frame
    ltf.word_wrap = True
    ltf.margin_left = Inches(0)
    p = ltf.paragraphs[0]
    configure_paragraph(p, "Existing Transit Systems", size=18, bold=True, color=LIGHT_ACCENT_TEAL, space_after=10)
    
    ex_list = [
        "Timetable PDF Sheets: Display rigid scheduling charts. Provide no actual status on traffic delays or cancellations.",
        "Official Tracking Apps (e.g. Gamyam): Provide coordinates mapping, but lack web PWA compatibility, working strictly on heavy app-store downloads.",
        "Hardware GPS Trackers: Require mounting proprietary tracking devices inside transit frames, costing high implementation investments."
    ]
    for item in ex_list:
        p = ltf.add_paragraph()
        configure_paragraph(p, item, size=13, color=LIGHT_TEXT_SECONDARY, space_after=10, level=1)

    # Right Box (Limitations)
    right_box = slide_4.shapes.add_textbox(Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.2))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    rtf.margin_left = Inches(0)
    p = rtf.paragraphs[0]
    configure_paragraph(p, "Limitations & Disadvantages", size=18, bold=True, color=LIGHT_ACCENT_TEAL, space_after=10)
    
    lim_list = [
        "No Active Driver Console: Standard drivers cannot feed GPS frames or log cabin occupant densities using a simple mobile UI.",
        "Heavy Bandwidth Needs: Advanced vector maps crash completely on rural routes under weak carrier signals.",
        "Hard API Key Dependencies: The passenger mapping screens fail to render if external mapping quotas expire.",
        "No Offline Simulation Mode: Testing route layouts requires live on-road vehicles with active internet syncs."
    ]
    for item in lim_list:
        p = rtf.add_paragraph()
        configure_paragraph(p, item, size=13, color=LIGHT_TEXT_SECONDARY, space_after=10, level=1)

    # ----------------------------------------------------
    # SLIDE 5: PROPOSED SYSTEM (KEY CAPABILITIES)
    # ----------------------------------------------------
    slide_5 = prs.slides.add_slide(blank_layout)
    apply_solid_background(slide_5, LIGHT_BG)
    add_slide_header(slide_5, "Proposed System: Key Capabilities")
    
    main_box = slide_5.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))
    mtf = main_box.text_frame
    mtf.word_wrap = True
    mtf.margin_left = Inches(0)
    
    p = mtf.paragraphs[0]
    configure_paragraph(p, "Dynamic Dual-Mode Synchronization Engine", size=16, bold=True, color=LIGHT_ACCENT_TEAL, space_after=2)
    p = mtf.add_paragraph()
    configure_paragraph(p, "Adapts dynamically: connects to Firebase Cloud Firestore for real-time live coordinate relays, or triggers client-side local database simulations offline.", size=14, color=LIGHT_TEXT_SECONDARY, space_after=12, level=1)
    
    p = mtf.add_paragraph()
    configure_paragraph(p, "Connection-Aware Mapping Framework", size=16, bold=True, color=LIGHT_ACCENT_TEAL, space_after=2)
    p = mtf.add_paragraph()
    configure_paragraph(p, "Continuously monitors cellular network status and automatically hot-swaps resource-heavy Google Maps for a lightweight Leaflet mapping environment.", size=14, color=LIGHT_TEXT_SECONDARY, space_after=12, level=1)
    
    p = mtf.add_paragraph()
    configure_paragraph(p, "High-Contrast Mobile Driver Duty Console", size=16, bold=True, color=LIGHT_ACCENT_TEAL, space_after=2)
    p = mtf.add_paragraph()
    configure_paragraph(p, "Responsive cabin dashboard built for operators. Streams coordinates and provides large touch buttons to increment and decrement bus load ratios.", size=14, color=LIGHT_TEXT_SECONDARY, space_after=12, level=1)
    
    p = mtf.add_paragraph()
    configure_paragraph(p, "Centralized Admin Control Room", size=16, bold=True, color=LIGHT_ACCENT_TEAL, space_after=2)
    p = mtf.add_paragraph()
    configure_paragraph(p, "Gives dispatchers full control: live map tracking, data seeding, incident management, and detailed CRUD editors for buses, drivers, stops, and schedules.", size=14, color=LIGHT_TEXT_SECONDARY, space_after=0, level=1)

    # ----------------------------------------------------
    # SLIDE 6: SYSTEM ARCHITECTURE (DESCRIPTIVE)
    # ----------------------------------------------------
    slide_6 = prs.slides.add_slide(blank_layout)
    apply_solid_background(slide_6, LIGHT_BG)
    add_slide_header(slide_6, "System Architecture: Layout & Blocks")
    
    main_box = slide_6.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))
    mtf = main_box.text_frame
    mtf.word_wrap = True
    mtf.margin_left = Inches(0)
    
    p = mtf.paragraphs[0]
    configure_paragraph(p, "1. User Roles (Client Interfaces)", size=16, bold=True, color=LIGHT_ACCENT_TEAL, space_after=2)
    p = mtf.add_paragraph()
    configure_paragraph(p, "Passenger UI (Route finder, live markers, stop sheets) | Driver UI (Telemetry streamer, logging console) | Admin UI (CRUD controls, global tracking)", size=13, color=LIGHT_TEXT_SECONDARY, space_after=10, level=1)
    
    p = mtf.add_paragraph()
    configure_paragraph(p, "2. React Application & Context Layer", size=16, bold=True, color=LIGHT_ACCENT_TEAL, space_after=2)
    p = mtf.add_paragraph()
    configure_paragraph(p, "AuthContext (Manages passenger, driver, and admin security states) | BusContext (Coordinates stops, routes, and active vehicles list) | MapContext (Coordinates bounds and map viewports)", size=13, color=LIGHT_TEXT_SECONDARY, space_after=10, level=1)
    
    p = mtf.add_paragraph()
    configure_paragraph(p, "3. Service Adapters & Streaming Libraries", size=16, bold=True, color=LIGHT_ACCENT_TEAL, space_after=2)
    p = mtf.add_paragraph()
    configure_paragraph(p, "mapsService.js (Coordinates Google Maps/Leaflet swaps) | gpsService.js (Controls HTML5 geolocation watcher) | realTransitData.js (Parses GTFS schedules) | busSimulator.js (Interpolates coordinates locally)", size=13, color=LIGHT_TEXT_SECONDARY, space_after=10, level=1)
    
    p = mtf.add_paragraph()
    configure_paragraph(p, "4. Persistent Database Layer", size=16, bold=True, color=LIGHT_ACCENT_TEAL, space_after=2)
    p = mtf.add_paragraph()
    configure_paragraph(p, "Cloud Firestore (Real-time NoSQL streaming websockets) | Browser LocalStorage (Local offline seeding fallback cache)", size=13, color=LIGHT_TEXT_SECONDARY, space_after=0, level=1)

    # ----------------------------------------------------
    # SLIDE 7: UML DIAGRAM (i) USE CASE DETAILS
    # ----------------------------------------------------
    slide_7 = prs.slides.add_slide(blank_layout)
    apply_solid_background(slide_7, LIGHT_BG)
    add_slide_header(slide_7, "UML Use Case Details")
    
    main_box = slide_7.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))
    mtf = main_box.text_frame
    mtf.word_wrap = True
    mtf.margin_left = Inches(0)
    
    p = mtf.paragraphs[0]
    configure_paragraph(p, "Passenger Portal Actor", size=16, bold=True, color=LIGHT_ACCENT_TEAL, space_after=2)
    p = mtf.add_paragraph()
    configure_paragraph(p, "Uses cases: [Search Transit Routes] | [Track Live Bus on Map] | [View Stop Timetables] | [Save Favorite Routes] | [Submit Incident Alert ticket]", size=13, color=LIGHT_TEXT_SECONDARY, space_after=12, level=1)
    
    p = mtf.add_paragraph()
    configure_paragraph(p, "Driver Duty Console Actor", size=16, bold=True, color=LIGHT_ACCENT_TEAL, space_after=2)
    p = mtf.add_paragraph()
    configure_paragraph(p, "Use cases: [Authenticate Driver Login] | [Start Journey Duty] | [End Journey Duty] | [Stream Live GPS Telemetry] | [Increment/Decrement Cabin Passenger Count]", size=13, color=LIGHT_TEXT_SECONDARY, space_after=12, level=1)
    
    p = mtf.add_paragraph()
    configure_paragraph(p, "Administrative Dashboard Actor", size=16, bold=True, color=LIGHT_ACCENT_TEAL, space_after=2)
    p = mtf.add_paragraph()
    configure_paragraph(p, "Use cases: [View Live Fleet Grid] | [Configure System Database Seeding] | [Manage CRUD Databases (Buses, Routes, Stops, Drivers)] | [Publish Scroll Announcements] | [Close Incident Tickets]", size=13, color=LIGHT_TEXT_SECONDARY, space_after=0, level=1)

    # ----------------------------------------------------
    # SLIDE 8: UML DIAGRAM (ii) CLASS STRUCTURE
    # ----------------------------------------------------
    slide_8 = prs.slides.add_slide(blank_layout)
    apply_solid_background(slide_8, LIGHT_BG)
    add_slide_header(slide_8, "UML Class Structure")
    
    main_box = slide_8.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))
    mtf = main_box.text_frame
    mtf.word_wrap = True
    mtf.margin_left = Inches(0)
    
    p = mtf.paragraphs[0]
    configure_paragraph(p, "Core Project Entities & Schema Mapping", size=16, bold=True, color=LIGHT_ACCENT_TEAL, space_after=10)
    
    classes = [
        "User Class: Contains user credentials, authentication profile (id, email, passwordHash, role: Passenger/Driver/Admin).",
        "Bus Class: Stores physical attributes (id, licensePlate, capacity) and operational tracking indicators (currentRouteId, driverId, status: active/inactive, passengerLoad, lastUpdated).",
        "Route Class: Maps individual routes (id, routeNumber, name, stopSequence: List of Stop IDs, durationMinutes, distanceKm).",
        "Stop Class: Maps coordinates to geographical stops (id, name, latitude, longitude, description, landmarks).",
        "Telemetry Class: Collects current streaming datasets (busId, lat, lng, speed, heading, timestamp). Methods: broadcastCoordinate(), calculateBearing(), calculateETA().",
        "Announcement Class: Stores global scroll tickers (id, message, priority: high/medium, timestamp, activeState)."
    ]
    for cls in classes:
        p = mtf.add_paragraph()
        configure_paragraph(p, cls, size=13, color=LIGHT_TEXT_SECONDARY, space_after=8, level=1)

    # ----------------------------------------------------
    # SLIDE 9: UML DIAGRAM (iii) ACTIVITY SEQUENCES (2 COLUMNS)
    # ----------------------------------------------------
    slide_9 = prs.slides.add_slide(blank_layout)
    apply_solid_background(slide_9, LIGHT_BG)
    add_slide_header(slide_9, "UML Activity Sequences")
    
    # Left Box (Driver)
    left_box = slide_9.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    ltf = left_box.text_frame
    ltf.word_wrap = True
    ltf.margin_left = Inches(0)
    p = ltf.paragraphs[0]
    configure_paragraph(p, "Driver Telemetry Stream Flow", size=18, bold=True, color=LIGHT_ACCENT_TEAL, space_after=10)
    
    drv_flow = [
        "1. Log in: Driver authenticates through console interface.",
        "2. Select Bus & Route: Driver assigns current vehicle assets.",
        "3. Start Duty: Geolocation watcher initializes browser tracking.",
        "4. Telemetry Loop (every 2-3 sec): GPS fetches coordinates -> streams speed/heading -> performs batch writes to Firestore.",
        "5. Cabin Count Check: Driver updates occupant volumes via counters.",
        "6. End Duty: Driver stops route -> coordinates streams terminate."
    ]
    for step in drv_flow:
        p = ltf.add_paragraph()
        configure_paragraph(p, step, size=13, color=LIGHT_TEXT_SECONDARY, space_after=8, level=1)

    # Right Box (Passenger)
    right_box = slide_9.shapes.add_textbox(Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.2))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    rtf.margin_left = Inches(0)
    p = rtf.paragraphs[0]
    configure_paragraph(p, "Passenger Map Tracking Flow", size=18, bold=True, color=LIGHT_ACCENT_TEAL, space_after=10)
    
    psg_flow = [
        "1. Open Commuter Portal: Passenger launches web app or PWA.",
        "2. Search Route: Commuter queries desired destination terminal.",
        "3. Map Bind: App mounts maps (Google Maps/Leaflet) to view.",
        "4. Real-time Listeners: BusContext monitors database for changes.",
        "5. Trigger Update: Active bus documents push new coordinates.",
        "6. Update Viewport: Marker relocates smoothly on map -> useETA hook updates arrival calculations instantly."
    ]
    for step in psg_flow:
        p = rtf.add_paragraph()
        configure_paragraph(p, step, size=13, color=LIGHT_TEXT_SECONDARY, space_after=8, level=1)

    # ----------------------------------------------------
    # SLIDE 10: UML DIAGRAM (iv) DEPLOYMENT LAYOUT
    # ----------------------------------------------------
    slide_10 = prs.slides.add_slide(blank_layout)
    apply_solid_background(slide_10, LIGHT_BG)
    add_slide_header(slide_10, "UML Deployment Layout")
    
    main_box = slide_10.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))
    mtf = main_box.text_frame
    mtf.word_wrap = True
    mtf.margin_left = Inches(0)
    
    p = mtf.paragraphs[0]
    configure_paragraph(p, "System Topology & Nodes Integration", size=16, bold=True, color=LIGHT_ACCENT_TEAL, space_after=10)
    
    nodes = [
        "Client Node (Driver Phone / Passenger Mobile Device): Runs browser instance (Chrome, Safari). Holds local service workers, HTML5 Geolocation API, and client cache LocalStorage.",
        "Application Web Server (Firebase Hosting Node): Serves bundled client-side React code chunks, CSS styles, and assets built by Vite.",
        "Cloud Storage Cluster Node (Firebase Cloud Firestore): Real-time synchronization cluster listening for updates. Pushes coordinate changes to subscribers via secure WebSocket protocols.",
        "Security Service Node (Firebase Authentication): Validates driver identities, checks role permissions, and distributes session authorization tokens (HTTPS protocols).",
        "Third-Party Mapping Service Providers (Google Maps Server & OpenStreetMap Tile Cluster): Feeds tile styles, mapping arrays, and geometry libraries directly to user browser containers."
    ]
    for node in nodes:
        p = mtf.add_paragraph()
        configure_paragraph(p, node, size=13, color=LIGHT_TEXT_SECONDARY, space_after=10, level=1)

    # ----------------------------------------------------
    # SLIDE 11: METHODOLOGY
    # ----------------------------------------------------
    slide_11 = prs.slides.add_slide(blank_layout)
    apply_solid_background(slide_11, LIGHT_BG)
    add_slide_header(slide_11, "Methodology & Development Phases")
    
    main_box = slide_11.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))
    mtf = main_box.text_frame
    mtf.word_wrap = True
    mtf.margin_left = Inches(0)
    
    p = mtf.paragraphs[0]
    configure_paragraph(p, "Phase 1: Transit Data Ingestion & Normalization", size=16, bold=True, color=LIGHT_ACCENT_TEAL, space_after=2)
    p = mtf.add_paragraph()
    configure_paragraph(p, "Parses raw TGSRTC GTFS schedules (.txt files containing routes, stops, schedules). Standardizes CSV matrices using custom parser scripts (realTransitData.js) and maps them into structured database rows.", size=13, color=LIGHT_TEXT_SECONDARY, space_after=12, level=1)
    
    p = mtf.add_paragraph()
    configure_paragraph(p, "Phase 2: Low-Overhead Telemetry Capture", size=16, bold=True, color=LIGHT_ACCENT_TEAL, space_after=2)
    p = mtf.add_paragraph()
    configure_paragraph(p, "Configures the Geolocation watcher to fire on coordinate drift. Batches telemetry commits to the database to keep write counts minimal while maintaining high location accuracy.", size=13, color=LIGHT_TEXT_SECONDARY, space_after=12, level=1)
    
    p = mtf.add_paragraph()
    configure_paragraph(p, "Phase 3: Connection-Quality Monitoring", size=16, bold=True, color=LIGHT_ACCENT_TEAL, space_after=2)
    p = mtf.add_paragraph()
    configure_paragraph(p, "Builds connection checkers (useConnectionQuality.js) to monitor connection status and ping times. Triggers an automatic switch from Google Maps to Leaflet maps on poor connection warnings.", size=13, color=LIGHT_TEXT_SECONDARY, space_after=12, level=1)
    
    p = mtf.add_paragraph()
    configure_paragraph(p, "Phase 4: Client-Side Simulation Engine", size=16, bold=True, color=LIGHT_ACCENT_TEAL, space_after=2)
    p = mtf.add_paragraph()
    configure_paragraph(p, "Integrates a client simulator (busSimulator.js) for testing in low-connectivity areas. Uses coordinate interpolation to simulate normal bus movements along routes offline.", size=13, color=LIGHT_TEXT_SECONDARY, space_after=0, level=1)

    # ----------------------------------------------------
    # SLIDE 12: IMPLEMENTATION & CODE DESIGN (2 COLUMNS)
    # ----------------------------------------------------
    slide_12 = prs.slides.add_slide(blank_layout)
    apply_solid_background(slide_12, LIGHT_BG)
    add_slide_header(slide_12, "Implementation & Code Design")
    
    # Left Box (Structure)
    left_box = slide_12.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    ltf = left_box.text_frame
    ltf.word_wrap = True
    ltf.margin_left = Inches(0)
    p = ltf.paragraphs[0]
    configure_paragraph(p, "Code Architecture & Services", size=18, bold=True, color=LIGHT_ACCENT_TEAL, space_after=10)
    
    impl_left = [
        "State Context Layers: AuthContext.jsx manages roles and sessions; BusContext.jsx handles routes and stops updates.",
        "Coordinate Services: gpsService.js streams coordinate packages; mapsService.js handles Google Maps and Leaflet bindings.",
        "Data Seeding Module: firebase.js sets up batch database writers; seedData.js coordinates mock route databases.",
        "Helper Methods: geoUtils.js calculates Haversine distances, bearings, and location offsets for transit lines."
    ]
    for item in impl_left:
        p = ltf.add_paragraph()
        configure_paragraph(p, item, size=13, color=LIGHT_TEXT_SECONDARY, space_after=10, level=1)

    # Right Box (Features)
    right_box = slide_12.shapes.add_textbox(Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.2))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    rtf.margin_left = Inches(0)
    p = rtf.paragraphs[0]
    configure_paragraph(p, "Implementation Progress & Status", size=18, bold=True, color=LIGHT_ACCENT_TEAL, space_after=10)
    
    impl_right = [
        "PWA Configuration: Offline page caching config is completed via the Vite PWA plugin, enabling offline application startup.",
        "Database Seeding Control: Admin controls can seed regional Hyderabad route stops with a single click.",
        "Real-Time Syncing: WebSocket pipelines for passenger map displays and driver telemetry trackers are operational.",
        "Responsive Designs: Dark-mode driver interfaces and mobile-friendly passenger views are completed using HSL color styling."
    ]
    for item in impl_right:
        p = rtf.add_paragraph()
        configure_paragraph(p, item, size=13, color=LIGHT_TEXT_SECONDARY, space_after=10, level=1)

    # ----------------------------------------------------
    # SLIDE 13: SCREENSHOTS PLACEHOLDERS
    # ----------------------------------------------------
    slide_13 = prs.slides.add_slide(blank_layout)
    apply_solid_background(slide_13, LIGHT_BG)
    add_slide_header(slide_13, "Interface Screenshots Layout")
    
    # Draw 4 styled boxes representing screenshots
    # Box 1 (Left Top): Passenger Portal Home
    sh1 = slide_13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.5), Inches(2.3))
    sh1.fill.solid()
    sh1.fill.fore_color.rgb = RGBColor(230, 235, 245)
    sh1.line.color.rgb = LIGHT_ACCENT_TEAL
    sh1.line.width = Pt(1.5)
    tf1 = sh1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    configure_paragraph(p, "PASSENGER HOME SCREEN", size=14, bold=True, color=LIGHT_TEXT_PRIMARY, space_after=4)
    p.alignment = PP_ALIGN.CENTER
    p = tf1.add_paragraph()
    configure_paragraph(p, "[Insert screenshot showing route selection autocomplete boxes, favorite route cards, and the bottom scrolling announcements ticker.]", size=11, color=LIGHT_TEXT_SECONDARY, space_after=0)
    p.alignment = PP_ALIGN.CENTER

    # Box 2 (Right Top): Passenger Live Map
    sh2 = slide_13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.6), Inches(5.5), Inches(2.3))
    sh2.fill.solid()
    sh2.fill.fore_color.rgb = RGBColor(230, 235, 245)
    sh2.line.color.rgb = LIGHT_ACCENT_TEAL
    sh2.line.width = Pt(1.5)
    tf2 = sh2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    configure_paragraph(p, "PASSENGER LIVE TRACKING MAP", size=14, bold=True, color=LIGHT_TEXT_PRIMARY, space_after=4)
    p.alignment = PP_ALIGN.CENTER
    p = tf2.add_paragraph()
    configure_paragraph(p, "[Insert screenshot showing the active map view, route paths, moving bus markers, and the real-time ETA information card.]", size=11, color=LIGHT_TEXT_SECONDARY, space_after=0)
    p.alignment = PP_ALIGN.CENTER

    # Box 3 (Left Bottom): Driver Console
    sh3 = slide_13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.3), Inches(5.5), Inches(2.3))
    sh3.fill.solid()
    sh3.fill.fore_color.rgb = RGBColor(230, 235, 245)
    sh3.line.color.rgb = LIGHT_ACCENT_TEAL
    sh3.line.width = Pt(1.5)
    tf3 = sh3.text_frame
    tf3.word_wrap = True
    p = tf3.paragraphs[0]
    configure_paragraph(p, "DRIVER TELEMETRY CONSOLE", size=14, bold=True, color=LIGHT_TEXT_PRIMARY, space_after=4)
    p.alignment = PP_ALIGN.CENTER
    p = tf3.add_paragraph()
    configure_paragraph(p, "[Insert screenshot of the dark-mode driver dashboard showing the active journey state, speed, compass heading, and passenger count buttons.]", size=11, color=LIGHT_TEXT_SECONDARY, space_after=0)
    p.alignment = PP_ALIGN.CENTER

    # Box 4 (Right Bottom): Admin Control Console
    sh4 = slide_13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(4.3), Inches(5.5), Inches(2.3))
    sh4.fill.solid()
    sh4.fill.fore_color.rgb = RGBColor(230, 235, 245)
    sh4.line.color.rgb = LIGHT_ACCENT_TEAL
    sh4.line.width = Pt(1.5)
    tf4 = sh4.text_frame
    tf4.word_wrap = True
    p = tf4.paragraphs[0]
    configure_paragraph(p, "ADMIN CONTROL ROOM DASHBOARD", size=14, bold=True, color=LIGHT_TEXT_PRIMARY, space_after=4)
    p.alignment = PP_ALIGN.CENTER
    p = tf4.add_paragraph()
    configure_paragraph(p, "[Insert screenshot of the admin control room dashboard showing the system metrics cards, database seeding panel, and CRUD editor grids.]", size=11, color=LIGHT_TEXT_SECONDARY, space_after=0)
    p.alignment = PP_ALIGN.CENTER

    # ----------------------------------------------------
    # SLIDE 14: FUTURE SCOPE & ENHANCEMENTS
    # ----------------------------------------------------
    slide_14 = prs.slides.add_slide(blank_layout)
    apply_solid_background(slide_14, LIGHT_BG)
    add_slide_header(slide_14, "Future Scope & Enhancements")
    
    main_box = slide_14.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.2))
    mtf = main_box.text_frame
    mtf.word_wrap = True
    mtf.margin_left = Inches(0)
    
    p = mtf.paragraphs[0]
    configure_paragraph(p, "Machine Learning-Based ETA Predictions", size=16, bold=True, color=LIGHT_ACCENT_TEAL, space_after=2)
    p = mtf.add_paragraph()
    configure_paragraph(p, "Incorporate machine learning regressors directly into the coordinates tracker, analyzing historical ride logs and time-of-day traffic to improve ETA forecasting.", size=13, color=LIGHT_TEXT_SECONDARY, space_after=12, level=1)
    
    p = mtf.add_paragraph()
    configure_paragraph(p, "Multi-Modal Public Transit Integration", size=16, bold=True, color=LIGHT_ACCENT_TEAL, space_after=2)
    p = mtf.add_paragraph()
    configure_paragraph(p, "Expand dataset ingestion modules to include regional MMTS train systems and the Hyderabad Metro Rail network, enabling seamless cross-transit trip planning.", size=13, color=LIGHT_TEXT_SECONDARY, space_after=12, level=1)
    
    p = mtf.add_paragraph()
    configure_paragraph(p, "Deviations & Disruption Alerts Dispatcher", size=16, bold=True, color=LIGHT_ACCENT_TEAL, space_after=2)
    p = mtf.add_paragraph()
    configure_paragraph(p, "Build algorithmic watchers to trigger alerts on the Admin Console if a driver deviates from the path coordinates, generating warnings automatically.", size=13, color=LIGHT_TEXT_SECONDARY, space_after=12, level=1)
    
    p = mtf.add_paragraph()
    configure_paragraph(p, "Unified QR-Code Ticketing Systems", size=16, bold=True, color=LIGHT_ACCENT_TEAL, space_after=2)
    p = mtf.add_paragraph()
    configure_paragraph(p, "Integrate QR payment widgets inside the Passenger Portal, enabling paperless ticketing verified by driver scanners at stop gates.", size=13, color=LIGHT_TEXT_SECONDARY, space_after=0, level=1)

    # ----------------------------------------------------
    # SLIDE 15: CONCLUSION (DARK SLATE THEME)
    # ----------------------------------------------------
    slide_15 = prs.slides.add_slide(blank_layout)
    apply_solid_background(slide_15, DARK_BG)
    
    # Large Thank You / Conclusion Header
    conclusion_box = slide_15.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.5))
    ctf = conclusion_box.text_frame
    ctf.word_wrap = True
    ctf.margin_left = Inches(0)
    
    p = ctf.paragraphs[0]
    configure_paragraph(p, "CONCLUSION", size=24, bold=True, color=ACCENT_TEAL, space_after=18)
    p = ctf.add_paragraph()
    configure_paragraph(p, "• Hardware-Independent Telemetry: CityBus successfully replaces expensive specialized GPS black boxes with driver-owned mobile devices, making transit tracking highly accessible.", size=15, color=DARK_TEXT_PRIMARY, space_after=12)
    p = ctf.add_paragraph()
    configure_paragraph(p, "• High Reliability: PWA capabilities and dynamic Leaflet fallbacks enable the application to work reliably, even in weak rural signals.", size=15, color=DARK_TEXT_PRIMARY, space_after=12)
    p = ctf.add_paragraph()
    configure_paragraph(p, "• Operations Simplified: Administrative control grids make fleet oversight and database updates straightforward for operators.", size=15, color=DARK_TEXT_PRIMARY, space_after=24)
    p = ctf.add_paragraph()
    configure_paragraph(p, "THANK YOU!", size=36, bold=True, color=ACCENT_TEAL, space_after=0)
    p.alignment = PP_ALIGN.CENTER
    
    # Save the presentation
    prs.save("CityBus_ProjectReview2.pptx")
    print("Successfully generated CityBus_ProjectReview2.pptx")

if __name__ == "__main__":
    create_presentation()
