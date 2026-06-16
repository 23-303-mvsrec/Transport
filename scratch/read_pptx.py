import sys
from pptx import Presentation

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    for i, slide in enumerate(prs.slides, 1):
        print(f"=== SLIDE {i} ===")
        # Try to get slide title
        if slide.shapes.title:
            print(f"Title: {slide.shapes.title.text.strip()}")
        else:
            print("Title: [No Title]")
            
        print("Content:")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        print(f"  - {text}")
        print("\n" + "="*40 + "\n")

if __name__ == "__main__":
    pptx_file = "D24_ProjectReview2 (1).pptx"
    extract_text_from_pptx(pptx_file)
